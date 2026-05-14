"""
Generate a professional PowerPoint about Backend Architecture for Senior Engineers.
Covers: SOLID in AI era, EDA, Kafka, Distributed Tracing, Microservices.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ──── Color Palette (Dark Professional) ────
BG_DARK    = RGBColor(0x0D, 0x11, 0x17)
BG_CARD    = RGBColor(0x16, 0x1B, 0x22)
ACCENT     = RGBColor(0x58, 0xA6, 0xFF)  # blue
GREEN      = RGBColor(0x3F, 0xB9, 0x50)
RED        = RGBColor(0xF8, 0x51, 0x49)
ORANGE     = RGBColor(0xF0, 0x88, 0x3E)
PURPLE     = RGBColor(0xBC, 0x8C, 0xFF)
YELLOW     = RGBColor(0xE3, 0xB3, 0x41)
WHITE      = RGBColor(0xE6, 0xED, 0xF3)
GRAY       = RGBColor(0x8B, 0x94, 0x9E)
DARK_GRAY  = RGBColor(0x30, 0x3D, 0x4F)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, left, top, width, height,
             font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=16, color=WHITE, bold=False, space_before=6, bullet=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Segoe UI'
    p.space_before = Pt(space_before)
    if bullet:
        p.level = 0
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_arrow(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Decorative accent line
add_rounded_rect(slide, 1, 1.5, 0.08, 2.5, ACCENT)

add_text(slide, 'BACKEND ARCHITECTURE', 1.5, 1.5, 10, 0.8,
         font_size=44, color=ACCENT, bold=True)
add_text(slide, 'Kiến trúc Backend cho Kỹ sư Cao cấp', 1.5, 2.3, 10, 0.6,
         font_size=28, color=WHITE, bold=True)

tf = add_text(slide, '', 1.5, 3.5, 8, 2, font_size=18, color=GRAY)
items = [
    '🏗️  SOLID Principles trong kỷ nguyên AI',
    '⚡  Event-Driven Architecture & Message Brokers',
    '🔍  Distributed Tracing & Observability',
    '🚀  Microservices Design Patterns',
]
for i, item in enumerate(items):
    add_para(tf, item, font_size=18, color=GRAY, space_before=10 if i > 0 else 0)

add_text(slide, 'Phạm Huỳnh Khánh', 1.5, 6.2, 4, 0.4, font_size=16, color=DARK_GRAY)
add_text(slide, '2026', 1.5, 6.6, 4, 0.4, font_size=14, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: Agenda / Overview
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 'NỘI DUNG TRÌNH BÀY', 0.8, 0.4, 10, 0.7,
         font_size=32, color=ACCENT, bold=True)

sections = [
    ('01', 'SOLID & AI', 'Tại sao kiến trúc phần mềm quan trọng\nhơn bao giờ hết trong kỷ nguyên AI', ACCENT),
    ('02', 'Event-Driven', 'Kafka, RabbitMQ và cách các\nservice giao tiếp bất đồng bộ', GREEN),
    ('03', 'Distributed Tracing', 'OpenTelemetry, Jaeger và cách\n"bắt bệnh" hệ thống phân tán', ORANGE),
    ('04', 'Microservices', 'Khi nào cần, khi nào không\nvà các pattern thiết yếu', PURPLE),
]

for i, (num, title, desc, color) in enumerate(sections):
    y = 1.5 + i * 1.4
    add_rounded_rect(slide, 0.8, y, 11.7, 1.15, BG_CARD, DARK_GRAY)
    add_text(slide, num, 1.2, y + 0.15, 0.8, 0.8, font_size=36, color=color, bold=True)
    add_text(slide, title, 2.2, y + 0.08, 3.5, 0.5, font_size=22, color=WHITE, bold=True)
    add_text(slide, desc, 5.5, y + 0.1, 6.5, 0.9, font_size=14, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: WHY SOLID MATTERS — Thợ xây vs Kiến trúc sư
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, '01  SOLID TRONG KỶ NGUYÊN AI', 0.8, 0.4, 10, 0.7,
         font_size=32, color=ACCENT, bold=True)
add_text(slide, 'Thợ xây  vs  Kiến trúc sư', 0.8, 1.1, 10, 0.5,
         font_size=20, color=GRAY)

# Left card — AI (Thợ xây)
add_rounded_rect(slide, 0.8, 1.8, 5.5, 4.5, BG_CARD, ACCENT)
add_text(slide, '🤖  AI = Thợ xây siêu việt', 1.2, 2.0, 4.8, 0.5,
         font_size=20, color=ACCENT, bold=True)
ai_items = [
    '✅ Viết code nhanh: 1800 dòng / 10 phút',
    '✅ Giải thuật toán phức tạp trong vài giây',
    '✅ Cấu hình boilerplate, CRUD tự động',
    '❌ Chỉ thấy context cục bộ',
    '❌ Không biết tầm nhìn 6 tháng của bạn',
    '❌ Dễ tạo "Spaghetti Code" nếu không chỉ đạo',
]
tf = add_text(slide, '', 1.2, 2.7, 4.8, 3.5, font_size=14, color=WHITE)
for item in ai_items:
    c = GREEN if item.startswith('✅') else RED
    add_para(tf, item, font_size=14, color=c, space_before=8)

# Right card — You (Kiến trúc sư)
add_rounded_rect(slide, 7.0, 1.8, 5.5, 4.5, BG_CARD, GREEN)
add_text(slide, '👨‍💻  BẠN = Kiến trúc sư', 7.4, 2.0, 4.8, 0.5,
         font_size=20, color=GREEN, bold=True)
you_items = [
    '🎯 Thiết kế hệ thống tổng thể',
    '🎯 Ra quyết định kiến trúc',
    '🎯 Chọn pattern phù hợp cho bài toán',
    '🎯 Review & validate output của AI',
    '🎯 Đảm bảo scalability dài hạn',
    '🎯 Prompt AI bằng "ngôn ngữ kiến trúc"',
]
tf = add_text(slide, '', 7.4, 2.7, 4.8, 3.5, font_size=14, color=WHITE)
for item in you_items:
    add_para(tf, item, font_size=14, color=WHITE, space_before=8)

# Bottom insight
add_rounded_rect(slide, 0.8, 6.5, 11.7, 0.7, DARK_GRAY)
add_text(slide, '💡  SOLID giúp bạn "chỉ đạo" AI chính xác — biến AI từ trợ lý thành công cụ mạnh mẽ', 1.2, 6.55, 11, 0.5,
         font_size=15, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: SOLID — Prompt Quality Comparison
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, '01  CHẤT LƯỢNG PROMPT = CHẤT LƯỢNG CODE', 0.8, 0.4, 12, 0.7,
         font_size=28, color=ACCENT, bold=True)

# Bad prompt
add_rounded_rect(slide, 0.8, 1.4, 5.5, 2.4, BG_CARD, RED)
add_text(slide, '❌  Prompt không có kiến thức kiến trúc', 1.2, 1.5, 5, 0.4,
         font_size=16, color=RED, bold=True)
add_text(slide, '"Viết cho tôi chức năng đăng nhập"', 1.2, 2.0, 5, 0.4,
         font_size=15, color=GRAY)
bad_results = [
    '→ AI viết 1 hàm khổng lồ 300 dòng',
    '→ Trộn lẫn: gọi API + validate + DB + render',
    '→ Không test được, không mở rộng được',
]
tf = add_text(slide, '', 1.2, 2.5, 5, 1.2, font_size=13, color=RED)
for item in bad_results:
    add_para(tf, item, font_size=13, color=RGBColor(0xFF, 0x99, 0x99), space_before=4)

# Good prompt
add_rounded_rect(slide, 7.0, 1.4, 5.5, 2.4, BG_CARD, GREEN)
add_text(slide, '✅  Prompt có tư duy kiến trúc', 7.4, 1.5, 5, 0.4,
         font_size=16, color=GREEN, bold=True)
add_text(slide, '"Tạo AuthService dùng Strategy Pattern\nđể switch giữa Google OAuth và Local Auth"', 7.4, 2.0, 5, 0.6,
         font_size=14, color=GRAY)
good_results = [
    '→ AI tạo interface + 2 implementation',
    '→ Mỗi class < 50 dòng, test riêng biệt',
    '→ Thêm Facebook Auth = thêm 1 class, 0 sửa code cũ',
]
tf = add_text(slide, '', 7.4, 2.7, 5, 1.2, font_size=13, color=GREEN)
for item in good_results:
    add_para(tf, item, font_size=13, color=RGBColor(0x99, 0xFF, 0x99), space_before=4)

# Blast Radius concept
add_text(slide, 'BLAST RADIUS — Bán kính sát thương khi sửa code', 0.8, 4.2, 12, 0.5,
         font_size=22, color=ORANGE, bold=True)

# Without SOLID
add_rounded_rect(slide, 0.8, 4.9, 5.5, 2.2, BG_CARD, RED)
add_text(slide, '💥  Không SOLID', 1.2, 5.0, 4, 0.4,
         font_size=16, color=RED, bold=True)
tf = add_text(slide, '', 1.2, 5.5, 5, 1.5, font_size=13, color=WHITE)
add_para(tf, 'AI sửa 1 endpoint trong file 500 dòng', font_size=13, color=WHITE, space_before=4)
add_para(tf, '→ Vô tình phá logic endpoint khác', font_size=13, color=RGBColor(0xFF, 0x99, 0x99), space_before=4)
add_para(tf, '→ Blast radius = toàn bộ module', font_size=13, color=RED, space_before=4, bold=True)

# With SOLID
add_rounded_rect(slide, 7.0, 4.9, 5.5, 2.2, BG_CARD, GREEN)
add_text(slide, '🛡️  Có SOLID', 7.4, 5.0, 4, 0.4,
         font_size=16, color=GREEN, bold=True)
tf = add_text(slide, '', 7.4, 5.5, 5, 1.5, font_size=13, color=WHITE)
add_para(tf, 'AI sửa 1 service class riêng biệt (50 dòng)', font_size=13, color=WHITE, space_before=4)
add_para(tf, '→ Các phần khác hoàn toàn an toàn', font_size=13, color=RGBColor(0x99, 0xFF, 0x99), space_before=4)
add_para(tf, '→ Blast radius = 0', font_size=13, color=GREEN, space_before=4, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: Event-Driven Architecture — Analogy
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, '02  EVENT-DRIVEN ARCHITECTURE (EDA)', 0.8, 0.4, 12, 0.7,
         font_size=32, color=GREEN, bold=True)
add_text(slide, 'Ví dụ Nhà hàng — Đồng bộ vs Bất đồng bộ', 0.8, 1.1, 10, 0.5,
         font_size=18, color=GRAY)

# Synchronous — bad
add_rounded_rect(slide, 0.8, 1.8, 5.5, 2.5, BG_CARD, RED)
add_text(slide, '❌  Giao tiếp trực tiếp (REST)', 1.2, 1.9, 5, 0.4,
         font_size=18, color=RED, bold=True)
tf = add_text(slide, '', 1.2, 2.4, 5, 1.8, font_size=14, color=GRAY)
add_para(tf, '🧑‍🍳 Khách gọi món', font_size=14, color=WHITE, space_before=6)
add_para(tf, '→ Nhân viên ĐỨNG CHỜ tại bếp', font_size=14, color=RGBColor(0xFF, 0x99, 0x99), space_before=4)
add_para(tf, '→ Không phục vụ ai khác được', font_size=14, color=RGBColor(0xFF, 0x99, 0x99), space_before=4)
add_para(tf, '→ Nghẽn cổ chai (Bottleneck)', font_size=14, color=RED, space_before=4, bold=True)

# Async — good
add_rounded_rect(slide, 7.0, 1.8, 5.5, 2.5, BG_CARD, GREEN)
add_text(slide, '✅  Event-Driven (Kafka/RabbitMQ)', 7.4, 1.9, 5, 0.4,
         font_size=18, color=GREEN, bold=True)
tf = add_text(slide, '', 7.4, 2.4, 5, 1.8, font_size=14, color=GRAY)
add_para(tf, '🧑‍🍳 Khách gọi món', font_size=14, color=WHITE, space_before=6)
add_para(tf, '→ Nhân viên ghi PHIẾU, thả vào rổ bếp', font_size=14, color=RGBColor(0x99, 0xFF, 0x99), space_before=4)
add_para(tf, '→ Đi phục vụ bàn khác ngay', font_size=14, color=RGBColor(0x99, 0xFF, 0x99), space_before=4)
add_para(tf, '→ Bất đồng bộ (Asynchronous)', font_size=14, color=GREEN, space_before=4, bold=True)

# 3 Components
add_text(slide, '3 THÀNH PHẦN CỐT LÕI', 0.8, 4.6, 10, 0.5,
         font_size=22, color=YELLOW, bold=True)

components = [
    ('PRODUCER', '🎯', 'Service tạo sự kiện\n"Đơn hàng #123 vừa tạo"', ACCENT),
    ('MESSAGE BROKER', '📬', 'Kafka / RabbitMQ\nLưu trữ & phân phối', ORANGE),
    ('CONSUMER', '⚡', 'Services lắng nghe\nInventory, Email, ...', GREEN),
]
for i, (title, emoji, desc, color) in enumerate(components):
    x = 0.8 + i * 4.2
    add_rounded_rect(slide, x, 5.2, 3.6, 2.0, BG_CARD, color)
    add_text(slide, emoji, x + 0.3, 5.3, 0.6, 0.5, font_size=28)
    add_text(slide, title, x + 1.0, 5.35, 2.3, 0.4, font_size=14, color=color, bold=True)
    add_text(slide, desc, x + 0.3, 5.9, 3.0, 1.0, font_size=12, color=GRAY)

# Arrows between components
add_arrow(slide, 4.5, 6.0, 0.4, 0.3, ACCENT)
add_arrow(slide, 8.7, 6.0, 0.4, 0.3, ORANGE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: EDA — Resilience (Chống đứt gãy)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, '02  SIÊU NĂNG LỰC: CHỐNG ĐỨT GÃY (Resilience)', 0.8, 0.4, 12, 0.7,
         font_size=28, color=GREEN, bold=True)

# Scenario: Email Service dies
add_text(slide, 'Kịch bản: Email Service bị sập', 0.8, 1.2, 10, 0.5,
         font_size=20, color=ORANGE, bold=True)

# Without EDA
add_rounded_rect(slide, 0.8, 1.9, 5.5, 2.5, BG_CARD, RED)
add_text(slide, '❌  Không có Message Broker', 1.2, 2.0, 5, 0.4,
         font_size=16, color=RED, bold=True)
tf = add_text(slide, '', 1.2, 2.5, 5, 1.6, font_size=13, color=WHITE)
add_para(tf, 'OrderService gọi EmailService', font_size=13, color=WHITE, space_before=6)
add_para(tf, '→ EmailService sập → trả lỗi', font_size=13, color=RGBColor(0xFF, 0x99, 0x99), space_before=4)
add_para(tf, '→ OrderService cũng sập theo', font_size=13, color=RGBColor(0xFF, 0x99, 0x99), space_before=4)
add_para(tf, '→ Khách không mua được hàng!', font_size=13, color=RED, space_before=4, bold=True)

# With EDA
add_rounded_rect(slide, 7.0, 1.9, 5.5, 2.5, BG_CARD, GREEN)
add_text(slide, '✅  Có Message Broker (Kafka)', 7.4, 2.0, 5, 0.4,
         font_size=16, color=GREEN, bold=True)
tf = add_text(slide, '', 7.4, 2.5, 5, 1.6, font_size=13, color=WHITE)
add_para(tf, 'OrderService ném event vào Kafka', font_size=13, color=WHITE, space_before=6)
add_para(tf, '→ Báo "Thành công!" cho khách', font_size=13, color=RGBColor(0x99, 0xFF, 0x99), space_before=4)
add_para(tf, '→ EmailService sập? Kafka giữ event an toàn', font_size=13, color=RGBColor(0x99, 0xFF, 0x99), space_before=4)
add_para(tf, '→ Sống lại → đọc tiếp → gửi email bù', font_size=13, color=GREEN, space_before=4, bold=True)

# ACK mechanism
add_text(slide, 'CƠ CHẾ XÁC NHẬN (ACK) — Tránh mất & lặp sự kiện', 0.8, 4.8, 12, 0.5,
         font_size=22, color=YELLOW, bold=True)

steps = [
    ('1', 'Broker gửi\nevent cho\nConsumer', ACCENT),
    ('2', 'Consumer\nxử lý\n(gửi email)', ORANGE),
    ('3', 'Consumer gửi\nACK về Broker\n"Đã xong!"', GREEN),
    ('4', 'Broker xóa\nevent khỏi\nhàng đợi', PURPLE),
]
for i, (num, desc, color) in enumerate(steps):
    x = 0.8 + i * 3.15
    add_rounded_rect(slide, x, 5.5, 2.6, 1.5, BG_CARD, color)
    add_circle(slide, x + 0.1, 5.6, 0.5, color)
    add_text(slide, num, x + 0.15, 5.6, 0.5, 0.5,
             font_size=20, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, desc, x + 0.7, 5.6, 1.7, 1.2, font_size=11, color=WHITE)

    if i < 3:
        add_arrow(slide, x + 2.7, 6.1, 0.35, 0.25, DARK_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: Distributed Tracing
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, '03  DISTRIBUTED TRACING', 0.8, 0.4, 12, 0.7,
         font_size=32, color=ORANGE, bold=True)
add_text(slide, '"Mò kim đáy bể" trong Microservices — TraceID là la bàn', 0.8, 1.1, 10, 0.5,
         font_size=18, color=GRAY)

# 3 Concepts
concepts = [
    ('Trace ID', '🛂', 'Hộ chiếu', 'Mã duy nhất đi theo request\ntừ đầu đến cuối toàn hệ thống', ACCENT),
    ('Span ID', '📍', 'Dấu mộc Visa', 'Mỗi service ghi lại:\nstart_time, end_time, SpanID', ORANGE),
    ('Context\nPropagation', '🔗', 'Truyền ngữ cảnh', 'REST → HTTP Header\ngRPC → Metadata\nKafka → Message Header', GREEN),
]
for i, (title, emoji, subtitle, desc, color) in enumerate(concepts):
    x = 0.8 + i * 4.2
    add_rounded_rect(slide, x, 1.8, 3.6, 2.8, BG_CARD, color)
    add_text(slide, emoji, x + 0.3, 1.9, 0.5, 0.5, font_size=30)
    add_text(slide, title, x + 1.0, 1.95, 2.3, 0.6, font_size=18, color=color, bold=True)
    add_text(slide, subtitle, x + 0.3, 2.6, 3, 0.4, font_size=13, color=YELLOW, bold=True)
    add_text(slide, desc, x + 0.3, 3.05, 3, 1.2, font_size=12, color=GRAY)

# Trace flow visualization
add_text(slide, 'VÍ DỤ LUỒNG TRACE', 0.8, 4.9, 10, 0.5,
         font_size=20, color=YELLOW, bold=True)

trace_flow = [
    ('API Gateway', 'TraceID: ABC\nSpanID: 1', ACCENT, '12ms'),
    ('Order Service', 'TraceID: ABC\nSpanID: 2, Parent: 1', GREEN, '45ms'),
    ('Kafka', 'TraceID: ABC\nin Message Header', ORANGE, '3ms'),
    ('Email Service', 'TraceID: ABC\nSpanID: 3, Parent: 2', PURPLE, '120ms'),
]
for i, (name, info, color, time) in enumerate(trace_flow):
    x = 0.5 + i * 3.2
    add_rounded_rect(slide, x, 5.5, 2.7, 1.6, BG_CARD, color)
    add_text(slide, name, x + 0.2, 5.55, 2.3, 0.35, font_size=13, color=color, bold=True)
    add_text(slide, info, x + 0.2, 5.95, 2.3, 0.7, font_size=10, color=GRAY)
    add_text(slide, f'⏱ {time}', x + 0.2, 6.7, 2.3, 0.3, font_size=11, color=YELLOW, bold=True)
    if i < 3:
        add_arrow(slide, x + 2.75, 6.1, 0.35, 0.25, DARK_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: OpenTelemetry & Tools
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, '03  OPENTELEMETRY — Bộ tiêu chuẩn vàng', 0.8, 0.4, 12, 0.7,
         font_size=28, color=ORANGE, bold=True)

add_rounded_rect(slide, 0.8, 1.3, 11.7, 1.5, BG_CARD, ORANGE)
add_text(slide, '💡  Trước đây: mỗi công cụ (Jaeger, Datadog, New Relic) yêu cầu code khác nhau → RẤT MỆT', 1.2, 1.4, 11, 0.4,
         font_size=14, color=RGBColor(0xFF, 0x99, 0x99))
add_text(slide, '✅  Bây giờ: OpenTelemetry (OTel) — 1 chuẩn chung, tự động hook vào HTTP / gRPC / Kafka / DB', 1.2, 1.9, 11, 0.4,
         font_size=14, color=RGBColor(0x99, 0xFF, 0x99))
add_text(slide, 'Cài thư viện → Auto-instrument → Gần như không cần code thêm cho tracing!', 1.2, 2.3, 11, 0.4,
         font_size=15, color=YELLOW, bold=True)

# Tools comparison
add_text(slide, 'CÔNG CỤ TRUY VẾT PHỔ BIẾN', 0.8, 3.2, 10, 0.5,
         font_size=22, color=ACCENT, bold=True)

tools = [
    ('Jaeger', 'Open-source\nCNCF Project\nMiễn phí', '🆓', GREEN),
    ('Zipkin', 'Open-source\nTwitter tạo ra\nNhẹ, đơn giản', '🪶', ACCENT),
    ('Datadog', 'Enterprise\nAll-in-one platform\nTrả phí', '💰', ORANGE),
    ('Grafana Tempo', 'Open-source\nTích hợp Grafana\nScalable', '📊', PURPLE),
]
for i, (name, desc, emoji, color) in enumerate(tools):
    x = 0.8 + i * 3.1
    add_rounded_rect(slide, x, 3.9, 2.6, 2.0, BG_CARD, color)
    add_text(slide, f'{emoji}  {name}', x + 0.2, 4.0, 2.2, 0.4, font_size=16, color=color, bold=True)
    add_text(slide, desc, x + 0.2, 4.5, 2.2, 1.2, font_size=12, color=GRAY)

# Node.js implementation note
add_rounded_rect(slide, 0.8, 6.2, 11.7, 1.0, DARK_GRAY)
add_text(slide, '🔧  Node.js: AsyncLocalStorage để truyền TraceID xuyên suốt request\n'
               '      → Batching: gom spans vào buffer RAM → bulk insert mỗi 3s (tránh nghẽn DB)', 1.2, 6.3, 11, 0.8,
         font_size=13, color=WHITE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: Microservices — When to use
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, '04  MICROSERVICES — KHI NÀO CẦN?', 0.8, 0.4, 12, 0.7,
         font_size=32, color=PURPLE, bold=True)

# Monolith vs Microservices
add_rounded_rect(slide, 0.8, 1.3, 5.5, 3.0, BG_CARD, ACCENT)
add_text(slide, '🏢  MONOLITHIC', 1.2, 1.4, 4.8, 0.4,
         font_size=20, color=ACCENT, bold=True)
add_text(slide, 'Phù hợp khi:', 1.2, 1.9, 4.8, 0.3,
         font_size=14, color=YELLOW, bold=True)
tf = add_text(slide, '', 1.2, 2.3, 4.8, 1.8, font_size=13, color=WHITE)
mono_items = [
    '✅ Team < 5 người',
    '✅ MVP / Startup giai đoạn đầu',
    '✅ Traffic < 10K req/s',
    '✅ Cần ship nhanh',
    '✅ Ít nghiệp vụ phức tạp',
]
for item in mono_items:
    add_para(tf, item, font_size=13, color=RGBColor(0x99, 0xFF, 0x99), space_before=4)

add_rounded_rect(slide, 7.0, 1.3, 5.5, 3.0, BG_CARD, PURPLE)
add_text(slide, '🔧  MICROSERVICES', 7.4, 1.4, 4.8, 0.4,
         font_size=20, color=PURPLE, bold=True)
add_text(slide, 'Phù hợp khi:', 7.4, 1.9, 4.8, 0.3,
         font_size=14, color=YELLOW, bold=True)
tf = add_text(slide, '', 7.4, 2.3, 4.8, 1.8, font_size=13, color=WHITE)
micro_items = [
    '✅ Team > 5 người, nhiều squad',
    '✅ Cần scale từng phần riêng biệt',
    '✅ Các domain nghiệp vụ độc lập',
    '✅ Polyglot (mỗi service ngôn ngữ khác)',
    '✅ Deploy độc lập, CI/CD riêng',
]
for item in micro_items:
    add_para(tf, item, font_size=13, color=RGBColor(0xCC, 0x99, 0xFF), space_before=4)

# 4 Pillars
add_text(slide, '4 TRỤ CỘT CỦA SENIOR BACKEND ENGINEER', 0.8, 4.6, 12, 0.5,
         font_size=22, color=YELLOW, bold=True)

pillars = [
    ('01', 'System\nArchitecture', 'Microservices\nEvent-Driven\nCAP Theorem', ACCENT),
    ('02', 'API Design\n& Integration', 'REST / gRPC\nGraphQL\nVersioning', GREEN),
    ('03', 'Performance\n& Operations', 'Caching, Indexing\nMonitoring\nIncident Handling', ORANGE),
    ('04', 'Technical\nLeadership', 'Code Review\nMentoring\nCross-functional', PURPLE),
]
for i, (num, title, desc, color) in enumerate(pillars):
    x = 0.8 + i * 3.1
    add_rounded_rect(slide, x, 5.2, 2.6, 2.1, BG_CARD, color)
    add_circle(slide, x + 0.2, 5.3, 0.5, color)
    add_text(slide, num, x + 0.25, 5.3, 0.5, 0.5,
             font_size=18, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.9, 5.3, 1.5, 0.7, font_size=12, color=color, bold=True)
    add_text(slide, desc, x + 0.2, 6.1, 2.2, 1.0, font_size=10, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: Kafka Delivery Semantics
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, '02  KAFKA — OFFSET & DELIVERY SEMANTICS', 0.8, 0.4, 12, 0.7,
         font_size=28, color=GREEN, bold=True)

# Offset concept
add_text(slide, 'Con trỏ Offset — "Bookmark" của Consumer', 0.8, 1.2, 10, 0.5,
         font_size=20, color=YELLOW, bold=True)

# Visual offset representation
offsets = [
    ('95', True), ('96', True), ('97', True), ('98', True), ('99', True),
    ('100', True), ('101', False), ('102', False), ('103', False),
]
for i, (num, processed) in enumerate(offsets):
    x = 0.8 + i * 1.3
    color = GREEN if processed else DARK_GRAY
    border = GREEN if processed else GRAY
    add_rounded_rect(slide, x, 1.9, 1.0, 0.6, color if processed else BG_CARD, border)
    add_text(slide, f'#{num}', x + 0.1, 1.95, 0.8, 0.4,
             font_size=14, color=BG_DARK if processed else GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# Pointer
add_text(slide, '▲', 7.3, 2.5, 0.5, 0.3, font_size=16, color=YELLOW, alignment=PP_ALIGN.CENTER)
add_text(slide, 'Consumer đang ở đây\n(Offset = 100)', 6.5, 2.8, 2, 0.6, font_size=11, color=YELLOW, alignment=PP_ALIGN.CENTER)

add_text(slide, '✅  Đã đọc & commit', 1.0, 2.7, 2, 0.3, font_size=11, color=GREEN)
add_text(slide, '⬜  Chưa đọc', 4.6, 2.7, 2, 0.3, font_size=11, color=GRAY)

# Delivery semantics table
add_text(slide, 'DELIVERY SEMANTICS — 3 chiến lược phân phối', 0.8, 3.6, 10, 0.5,
         font_size=20, color=ACCENT, bold=True)

semantics = [
    ('At-most-once', 'Gửi 1 lần, không retry\nCó thể MẤT event', '⚡ Nhanh nhất\n❌ Rủi ro mất dữ liệu', RED),
    ('At-least-once', 'Gửi lại nếu chưa ACK\nCó thể GỬI LẶP', '✅ Không mất event\n⚠️ Cần Idempotency', ORANGE),
    ('Exactly-once', 'Kafka Transactions\nĐảm bảo đúng 1 lần', '🏆 An toàn nhất\n❌ Chậm nhất, phức tạp', GREEN),
]
for i, (name, desc, tradeoff, color) in enumerate(semantics):
    y = 4.2 + i * 1.05
    add_rounded_rect(slide, 0.8, y, 11.7, 0.85, BG_CARD, color)
    add_text(slide, name, 1.2, y + 0.15, 2.5, 0.5, font_size=16, color=color, bold=True)
    add_text(slide, desc, 3.8, y + 0.05, 4, 0.7, font_size=11, color=GRAY)
    add_text(slide, tradeoff, 8.2, y + 0.05, 4, 0.7, font_size=11, color=WHITE)

# Idempotency callout
add_rounded_rect(slide, 0.8, 7.4, 11.7, 0.0, DARK_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: Summary & Roadmap
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 'TÓM TẮT & LỘ TRÌNH HỌC TẬP', 0.8, 0.4, 12, 0.7,
         font_size=32, color=ACCENT, bold=True)

# Key takeaways
add_text(slide, 'KEY TAKEAWAYS', 0.8, 1.2, 10, 0.5,
         font_size=20, color=YELLOW, bold=True)

takeaways = [
    ('SOLID + AI', 'Kiến trúc phần mềm giúp "chỉ đạo" AI chính xác, giảm blast radius', ACCENT),
    ('EDA + Kafka', 'Hệ thống bất đồng bộ, chống đứt gãy, scale horizontal dễ dàng', GREEN),
    ('Tracing', 'TraceID + SpanID + Context Propagation = la bàn debug hệ thống phân tán', ORANGE),
    ('Microservices', 'Không phải silver bullet — chỉ dùng khi thực sự cần scale team/traffic', PURPLE),
]
for i, (title, desc, color) in enumerate(takeaways):
    y = 1.8 + i * 0.75
    add_rounded_rect(slide, 0.8, y, 11.7, 0.6, BG_CARD, color)
    add_text(slide, title, 1.2, y + 0.07, 2.5, 0.4, font_size=14, color=color, bold=True)
    add_text(slide, desc, 3.8, y + 0.07, 8.4, 0.4, font_size=13, color=GRAY)

# Learning roadmap
add_text(slide, 'LỘ TRÌNH THỰC TẾ', 0.8, 5.0, 10, 0.5,
         font_size=20, color=YELLOW, bold=True)

roadmap = [
    ('Junior', 'SOLID\nREST API\nSQL/NoSQL', ACCENT),
    ('Mid', 'Design Patterns\nDocker/CI-CD\nCaching/Redis', GREEN),
    ('Senior', 'Microservices\nKafka/EDA\nSystem Design', ORANGE),
    ('Architect', 'Distributed Tracing\nCAP/SAGA\nTech Leadership', PURPLE),
]
for i, (level, skills, color) in enumerate(roadmap):
    x = 0.8 + i * 3.1
    add_rounded_rect(slide, x, 5.6, 2.6, 1.6, BG_CARD, color)
    add_text(slide, level, x + 0.2, 5.65, 2.2, 0.35, font_size=16, color=color, bold=True)
    add_text(slide, skills, x + 0.2, 6.05, 2.2, 1.0, font_size=11, color=GRAY)
    if i < 3:
        add_arrow(slide, x + 2.7, 6.2, 0.3, 0.2, DARK_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12: Thank you
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_rounded_rect(slide, 1, 2.5, 0.08, 2.5, ACCENT)
add_text(slide, 'CẢM ƠN!', 1.5, 2.5, 10, 1,
         font_size=52, color=ACCENT, bold=True)
add_text(slide, 'Q & A', 1.5, 3.5, 10, 0.7,
         font_size=28, color=WHITE)

tf = add_text(slide, '', 1.5, 4.5, 8, 1.5, font_size=16, color=GRAY)
add_para(tf, '📧  phamhuynhkhanh@email.com', font_size=16, color=GRAY, space_before=8)
add_para(tf, '💻  github.com/phamhuynhkhanh', font_size=16, color=GRAY, space_before=8)

# Save
output_path = os.path.join(os.path.dirname(__file__), 'Backend_Architecture_Presentation.pptx')
prs.save(output_path)
print(f'Done! Saved to: {output_path}')
